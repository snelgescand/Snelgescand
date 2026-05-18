#!/usr/bin/env python3
"""
generate_pptx.py — PowerPoint-generator voor Sportief Opgewekt.

Leest:
  - Pad naar TEMPLATE.pptx (eerste argument)
  - JSON op stdin met de project-data + berekend resultaat

Schrijft:
  - .pptx bytes op stdout

Aangeroepen vanuit Node:
    spawn('python3', ['generate_pptx.py', '/path/to/template.pptx'], { stdio: ['pipe', 'pipe', 'inherit'] })

Placeholder-syntax in de template (in elke tekstrun, ook in tabelcellen en
notities-pagina):

    {{club.naam}}                          → string
    {{rollup.nettoInvestering | euro}}     → € 128.556
    {{rollup.totaleCo2BesparingKg | ton}}  → 26.9 ton/jaar
    {{rollup.gemiddeldeTerugverdientijdJaren | jaren}}  → 8.4 jaar
    {{maatregel.dakisolatie.besparingPerJaar | euro}}   → € 416

Filters: euro, ton, kwh, m3, jaren, pct, datum, geen-filter (raw)

Slides verbergen op basis van condities:
    Slide-notitie bevat:  ALLEEN_ALS: rollup.aansluitwaardeVoldoende == false
    → slide blijft staan (verbergen werkt niet in python-pptx; sectie blijft)

Diagrammen worden NIET dynamisch ingevuld in deze MVP. Voor sprint 7+ kan
matplotlib of een PNG-injectie toegevoegd worden.

Vereist:
    pip install python-pptx==0.6.23
"""

import json
import re
import sys
from copy import deepcopy
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401  (gereserveerd voor latere styling)
except ImportError:
    print("ERROR: python-pptx niet geïnstalleerd. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(2)


# ============================================================================
# Filters
# ============================================================================

def filter_euro(value: Any) -> str:
    if value is None or value == "":
        return "—"
    try:
        n = float(value)
    except (TypeError, ValueError):
        return str(value)
    # NL-notatie: € 128.556 (geen decimalen voor grote bedragen)
    return "€ {:,.0f}".format(n).replace(",", ".")


def filter_ton(value: Any) -> str:
    if value is None:
        return "—"
    try:
        n = float(value) / 1000.0
    except (TypeError, ValueError):
        return str(value)
    return "{:,.1f} ton".format(n).replace(",", ".")


def filter_kwh(value: Any) -> str:
    if value is None:
        return "—"
    try:
        n = float(value)
    except (TypeError, ValueError):
        return str(value)
    return "{:,.0f} kWh".format(n).replace(",", ".")


def filter_m3(value: Any) -> str:
    if value is None:
        return "—"
    try:
        n = float(value)
    except (TypeError, ValueError):
        return str(value)
    return "{:,.0f} m³".format(n).replace(",", ".")


def filter_jaren(value: Any) -> str:
    if value is None:
        return "—"
    try:
        n = float(value)
    except (TypeError, ValueError):
        return str(value)
    if n != n or n == float("inf"):  # NaN of Infinity
        return "n.v.t."
    if n >= 100:
        return ">100 jaar"
    return "{:.1f} jaar".format(n).replace(".", ",")


def filter_pct(value: Any) -> str:
    if value is None:
        return "—"
    try:
        n = float(value) * 100
    except (TypeError, ValueError):
        return str(value)
    return "{:.0f}%".format(n)


def filter_int(value: Any) -> str:
    if value is None:
        return "—"
    try:
        return "{:,.0f}".format(float(value)).replace(",", ".")
    except (TypeError, ValueError):
        return str(value)


FILTERS = {
    "euro": filter_euro,
    "ton": filter_ton,
    "kwh": filter_kwh,
    "m3": filter_m3,
    "jaren": filter_jaren,
    "pct": filter_pct,
    "int": filter_int,
    "raw": lambda v: "" if v is None else str(v),
}


# ============================================================================
# Placeholder-resolutie
# ============================================================================

PLACEHOLDER_RE = re.compile(r"\{\{\s*([a-zA-Z0-9_.\[\]\-]+)(?:\s*\|\s*([a-zA-Z0-9_]+))?\s*\}\}")


def resolve_path(data: dict, path: str) -> Any:
    """Resolve a dotted path like 'rollup.nettoInvestering' against data dict."""
    parts = path.split(".")
    current: Any = data
    for part in parts:
        if current is None:
            return None
        if isinstance(current, dict):
            current = current.get(part)
        elif isinstance(current, list):
            try:
                idx = int(part)
                current = current[idx] if 0 <= idx < len(current) else None
            except (ValueError, IndexError):
                return None
        else:
            attr = getattr(current, part, None)
            current = attr
    return current


def substitute_text(text: str, data: dict) -> str:
    """Replace all {{path|filter}} placeholders in a string."""
    def repl(match: re.Match) -> str:
        path = match.group(1)
        filter_name = match.group(2) or "raw"
        value = resolve_path(data, path)
        filter_fn = FILTERS.get(filter_name)
        if filter_fn is None:
            return f"[onbekend filter: {filter_name}]"
        return filter_fn(value)
    return PLACEHOLDER_RE.sub(repl, text)


# ============================================================================
# Slide-traversal: vervang alle placeholders in tekstrun's
# ============================================================================

def process_text_frame(text_frame, data: dict) -> None:
    """
    Vervang placeholders in alle tekstrun's.

    Belangrijk: python-pptx splitst een placeholder soms over meerdere runs
    als de gebruiker midden in {{...}} een spatie of stijl heeft gewijzigd.
    Onze strategie: voeg run-tekst per paragraph samen, vervang in de
    samengevoegde string, en zet het resultaat terug in de eerste run.
    """
    if text_frame is None:
        return
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            continue
        full_text = "".join(run.text for run in paragraph.runs)
        if "{{" not in full_text:
            continue
        new_text = substitute_text(full_text, data)
        if new_text == full_text:
            continue
        # Zet samengevoegde tekst in eerste run, leeg de rest.
        # Stijl van eerste run blijft behouden.
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""


def process_shape(shape, data: dict) -> None:
    """Recursief alle vormen + group-shapes verwerken."""
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, data)
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            process_shape(sub, data)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, data)


def process_slide(slide, data: dict) -> None:
    for shape in slide.shapes:
        process_shape(shape, data)
    # Speaker notes
    if slide.has_notes_slide:
        process_text_frame(slide.notes_slide.notes_text_frame, data)


# ============================================================================
# Main
# ============================================================================

def main() -> int:
    if len(sys.argv) < 2:
        print("Gebruik: generate_pptx.py <template.pptx>", file=sys.stderr)
        return 2

    template_path = Path(sys.argv[1])
    if not template_path.exists():
        print(f"Template niet gevonden: {template_path}", file=sys.stderr)
        return 2

    # Project + resultaat JSON van stdin
    try:
        raw = sys.stdin.read()
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        print(f"Ongeldige JSON op stdin: {exc}", file=sys.stderr)
        return 2

    # Maak per-maatregel lookup gemakkelijker: data['maatregel']['<id>']
    if "perMaatregel" in data and "maatregel" not in data:
        data["maatregel"] = deepcopy(data["perMaatregel"])

    # Open template, vervang, schrijf naar stdout
    presentation = Presentation(str(template_path))
    for slide in presentation.slides:
        process_slide(slide, data)

    presentation.save(sys.stdout.buffer)
    return 0


if __name__ == "__main__":
    sys.exit(main())
