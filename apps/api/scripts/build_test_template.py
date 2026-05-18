#!/usr/bin/env python3
"""
build_test_template.py — bouwt een minimale TEMPLATE_v2.pptx met placeholders.

Dit is NIET de uiteindelijke huisstijl-template (die maakt Bart zelf op
basis van de bestaande 86-slide template). Dit is een test-template om te
bewijzen dat de pipeline werkt.

Gebruik:
    python3 build_test_template.py  →  schrijft TEMPLATE_v2.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "TEMPLATE_v2.pptx"

# Sportief Opgewekt-achtige kleuren — VERVANG met echte huisstijl
PRIMARY = RGBColor(0x16, 0xA3, 0x4A)  # groen
DARK = RGBColor(0x14, 0x53, 0x2D)
GRAY = RGBColor(0x64, 0x74, 0x8B)


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=DARK):
    """Voeg een textbox toe met opgegeven tekst en styling."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Voorblad
    s1 = prs.slides.add_slide(blank)
    add_text(s1, 0.5, 0.5, 12, 0.5, "Verduurzamingsplan", size=14, color=GRAY)
    add_text(s1, 0.5, 1.2, 12, 1, "{{context.club.naam}}", size=44, bold=True, color=PRIMARY)
    add_text(s1, 0.5, 2.5, 12, 0.5, "Opgesteld voor: {{context.club.naam}}", size=18, color=DARK)
    add_text(s1, 0.5, 6.5, 12, 0.5, "Sportief Opgewekt · Snelgescand.nl", size=12, color=GRAY)

    # Slide 2 — Uitgangspunten
    s2 = prs.slides.add_slide(blank)
    add_text(s2, 0.5, 0.4, 12, 0.7, "Uitgangspunten", size=32, bold=True, color=PRIMARY)
    add_text(s2, 0.5, 1.5, 6, 0.5, "Bouwjaar clubhuis:", size=14, color=GRAY)
    add_text(s2, 4, 1.5, 6, 0.5, "{{context.gebouw.bouwjaar}}", size=14, bold=True)
    add_text(s2, 0.5, 2.1, 6, 0.5, "Bruto vloeroppervlak:", size=14, color=GRAY)
    add_text(s2, 4, 2.1, 6, 0.5, "{{context.gebouw.bvoTotaalM2 | int}} m²", size=14, bold=True)
    add_text(s2, 0.5, 2.7, 6, 0.5, "Gasverbruik per jaar:", size=14, color=GRAY)
    add_text(s2, 4, 2.7, 6, 0.5, "{{context.energie.gasverbruikM3 | m3}}", size=14, bold=True)
    add_text(s2, 0.5, 3.3, 6, 0.5, "Stroomverbruik per jaar:", size=14, color=GRAY)
    add_text(s2, 4, 3.3, 6, 0.5, "{{context.energie.stroomverbruikTotaalKwh | kwh}}", size=14, bold=True)

    # Slide 3 — Voor de penningmeester
    s3 = prs.slides.add_slide(blank)
    add_text(s3, 0.5, 0.4, 12, 0.7, "Voor de penningmeester", size=32, bold=True, color=PRIMARY)
    add_text(s3, 0.5, 1.5, 6, 0.5, "Bruto investering:", size=16, color=GRAY)
    add_text(s3, 6, 1.5, 6, 0.5, "{{rollup.totaleInvestering | euro}}", size=18, bold=True)
    add_text(s3, 0.5, 2.2, 6, 0.5, "Totale subsidies:", size=16, color=GRAY)
    add_text(s3, 6, 2.2, 6, 0.5, "{{rollup.totaleSubsidie | euro}}", size=18, bold=True)
    add_text(s3, 0.5, 2.9, 6, 0.5, "Netto investering:", size=16, color=GRAY)
    add_text(s3, 6, 2.9, 6, 0.5, "{{rollup.nettoInvestering | euro}}", size=20, bold=True, color=PRIMARY)
    add_text(s3, 0.5, 3.8, 6, 0.5, "Besparing per jaar:", size=16, color=GRAY)
    add_text(s3, 6, 3.8, 6, 0.5, "{{rollup.totaleBesparingPerJaar | euro}}", size=18, bold=True)
    add_text(s3, 0.5, 4.5, 6, 0.5, "Gemiddelde terugverdientijd:", size=16, color=GRAY)
    add_text(s3, 6, 4.5, 6, 0.5, "{{rollup.gemiddeldeTerugverdientijdJaren | jaren}}", size=18, bold=True)
    add_text(s3, 0.5, 5.2, 6, 0.5, "CO₂-besparing:", size=16, color=GRAY)
    add_text(s3, 6, 5.2, 6, 0.5, "{{rollup.totaleCo2BesparingKg | ton}}/jaar", size=18, bold=True)

    # Slide 4 — Voorbeeld per-maatregel
    s4 = prs.slides.add_slide(blank)
    add_text(s4, 0.5, 0.4, 12, 0.7, "Maatregel uitgelicht: dakisolatie", size=28, bold=True, color=PRIMARY)
    add_text(s4, 0.5, 1.5, 5, 0.5, "Bruto investering:", size=14, color=GRAY)
    add_text(s4, 5, 1.5, 4, 0.5, "{{maatregel.dakisolatie.brutoInvestering | euro}}", size=14, bold=True)
    add_text(s4, 0.5, 2.1, 5, 0.5, "Subsidie:", size=14, color=GRAY)
    add_text(s4, 5, 2.1, 4, 0.5, "{{maatregel.dakisolatie.totaleSubsidie | euro}}", size=14, bold=True)
    add_text(s4, 0.5, 2.7, 5, 0.5, "Besparing per jaar:", size=14, color=GRAY)
    add_text(s4, 5, 2.7, 4, 0.5, "{{maatregel.dakisolatie.besparingPerJaar | euro}}", size=14, bold=True)
    add_text(s4, 0.5, 3.3, 5, 0.5, "Terugverdientijd:", size=14, color=GRAY)
    add_text(s4, 5, 3.3, 4, 0.5, "{{maatregel.dakisolatie.terugverdientijdJaren | jaren}}", size=14, bold=True)

    prs.save(str(OUT))
    print(f"✓ Template geschreven: {OUT}")


if __name__ == "__main__":
    main()
